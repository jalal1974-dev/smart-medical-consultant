"""
ULTIMATE SMART MEDICAL CONSULTANT - BACKEND API
================================================
Professional medical consultation platform with Claude AI
Generates perfect Arabic/English PPTX presentations and infographics
Matching your exact brand design!

Author: Built for Smart Medical Consultant (مستشارك الطبي الذكي)
Version: 2.0 Ultimate Edition
"""

from fastapi import FastAPI, HTTPException, UploadFile, File
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel, Field
from typing import Optional, List, Literal, Dict
import anthropic
import os
from io import BytesIO
import base64
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import logging
from datetime import datetime

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Initialize FastAPI
app = FastAPI(
    title="Smart Medical Consultant API",
    description="Ultimate medical consultation platform with AI-powered analysis",
    version="2.0.0"
)

# CORS Configuration
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],  # TODO: Restrict in production
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Brand Colors (from your PPTX analysis)
class BrandColors:
    CYAN_HEADER = RGBColor(6, 182, 212)      # #06B6D4
    CYAN_ACCENT = RGBColor(8, 145, 178)       # #0891B2
    GREEN_SUCCESS = RGBColor(16, 185, 129)    # #10B981
    GREEN_LIGHT = RGBColor(209, 250, 229)     # #D1FAE5
    GREEN_DARK = RGBColor(6, 95, 70)          # #065F46
    TEXT_PRIMARY = RGBColor(22, 78, 99)       # #164E63
    TEXT_SECONDARY = RGBColor(100, 116, 139)  # #64748B
    WHITE = RGBColor(255, 255, 255)           # #FFFFFF
    BACKGROUND = RGBColor(248, 250, 252)      # #F8FAFC

# API Configuration
ANTHROPIC_API_KEY = os.getenv("ANTHROPIC_API_KEY")
if not ANTHROPIC_API_KEY:
    logger.warning("⚠️ ANTHROPIC_API_KEY not set - API calls will fail!")
else:
    client = anthropic.Anthropic(api_key=ANTHROPIC_API_KEY)

# ==================== MODELS ====================

class PatientData(BaseModel):
    """Patient information model"""
    name: str = Field(..., description="Patient name")
    symptoms: List[str] = Field(..., description="List of symptoms")
    medical_history: Optional[str] = Field(None, description="Medical history")
    diagnosis: str = Field(..., description="Medical diagnosis")
    recommendations: List[str] = Field(..., description="Treatment recommendations")
    tests: List[str] = Field(..., description="Recommended medical tests")
    urgency: Optional[Literal["low", "medium", "high"]] = Field("medium", description="Urgency level")

class InfographicRequest(BaseModel):
    """Infographic generation request"""
    patient_data: PatientData
    language: Literal["en", "ar"] = Field("ar", description="Content language")
    custom_prompt: Optional[str] = Field(None, description="Custom generation instructions")

class SlideRequest(BaseModel):
    """Slide deck generation request"""
    patient_data: PatientData
    consultation_id: int
    language: Literal["en", "ar"] = Field("ar", description="Content language")

class GenerateResponse(BaseModel):
    """API response model"""
    success: bool
    data: Optional[str] = Field(None, description="Base64 encoded file")
    file_type: Optional[str] = Field(None, description="File MIME type")
    size_bytes: Optional[int] = Field(None, description="File size in bytes")
    error: Optional[str] = Field(None, description="Error message if failed")

# ==================== HEALTH CHECK ====================

@app.get("/", tags=["Health"])
async def root():
    """API root endpoint"""
    return {
        "service": "Smart Medical Consultant API",
        "version": "2.0.0 Ultimate Edition",
        "status": "operational",
        "tagline": "مستشارك الطبي الذكي | Premium Healthcare Technology"
    }

@app.get("/health", tags=["Health"])
async def health_check():
    """Health check endpoint"""
    return {
        "status": "healthy",
        "timestamp": datetime.utcnow().isoformat(),
        "anthropic_configured": bool(ANTHROPIC_API_KEY)
    }

# ==================== INFOGRAPHIC GENERATION ====================

@app.post("/generate/infographic", response_model=GenerateResponse, tags=["Generation"])
async def generate_infographic(request: InfographicRequest):
    """Generate medical infographic as SVG"""
    try:
        logger.info(f"🎨 Generating infographic - Language: {request.language}")
        
        # Generate SVG using Claude
        svg_content = await generate_infographic_svg(
            patient_data=request.patient_data.dict(),
            language=request.language,
            custom_prompt=request.custom_prompt
        )
        
        svg_base64 = base64.b64encode(svg_content.encode('utf-8')).decode('utf-8')
        logger.info("✅ Infographic generated successfully")
        
        return GenerateResponse(
            success=True,
            data=svg_base64,
            file_type="image/svg+xml",
            size_bytes=len(svg_content)
        )
        
    except Exception as e:
        logger.error(f"❌ Infographic generation failed: {str(e)}")
        return GenerateResponse(success=False, error=str(e))

# ==================== SLIDE GENERATION ====================

@app.post("/generate/slides", response_model=GenerateResponse, tags=["Generation"])
async def generate_slides(request: SlideRequest):
    """Generate PowerPoint presentation"""
    try:
        logger.info(f"📊 Generating PPTX - Language: {request.language}, ID: {request.consultation_id}")
        
        prs = create_ultimate_medical_presentation(
            patient_data=request.patient_data,
            consultation_id=request.consultation_id,
            language=request.language
        )
        
        pptx_bytes = BytesIO()
        prs.save(pptx_bytes)
        pptx_bytes.seek(0)
        
        file_data = pptx_bytes.read()
        pptx_base64 = base64.b64encode(file_data).decode('utf-8')
        
        logger.info(f"✅ PPTX generated - Size: {len(file_data)} bytes")
        
        return GenerateResponse(
            success=True,
            data=pptx_base64,
            file_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            size_bytes=len(file_data)
        )
        
    except Exception as e:
        logger.error(f"❌ PPTX generation failed: {str(e)}")
        return GenerateResponse(success=False, error=str(e))

# ==================== INFOGRAPHIC GENERATOR ====================

async def generate_infographic_svg(patient_data: dict, language: str, custom_prompt: Optional[str] = None) -> str:
    """Generate medical infographic SVG using Claude"""
    
    is_arabic = language == "ar"
    
    if is_arabic:
        system_prompt = """أنت مصمم جرافيك طبي محترف متخصص في التصميم الطبي.

قم بإنشاء رسم بياني طبي (إنفوجرافيك) احترافي بصيغة SVG.

المواصفات الفنية:
- الأبعاد: 800 بكسل عرض × 1200 بكسل ارتفاع
- الألوان: الأزرق السماوي #06B6D4، الأخضر #10B981، الأبيض #FFFFFF
- النص: عربي فقط بدون أي كلمة إنجليزية
- الاتجاه: من اليمين لليسار (RTL)

أنشئ تصميماً طبياً احترافياً يتضمن:
1. عنوان بارز في الأعلى
2. بيانات المريض
3. الأعراض مع أيقونات
4. التشخيص في إطار ملون
5. التوصيات العلاجية
6. الفحوصات المطلوبة

أرجع كود SVG فقط بدون أي شرح."""

        patient_name = patient_data.get('name', 'المريض')
        symptoms_text = "، ".join(patient_data.get('symptoms', [])[:5])
        diagnosis = patient_data.get('diagnosis', 'يتطلب تقييم طبي')
        recommendations = patient_data.get('recommendations', [])[:4]
        tests = patient_data.get('tests', [])[:4]
        
        user_prompt = f"""البيانات الطبية:

المريض: {patient_name}
الأعراض: {symptoms_text}
التشخيص: {diagnosis}
التوصيات: {', '.join(recommendations)}
الفحوصات: {', '.join(tests)}

{f"ملاحظات إضافية: {custom_prompt}" if custom_prompt else ""}

أنشئ إنفوجرافيك طبي احترافي باللغة العربية فقط."""
    
    else:
        system_prompt = """You are a professional medical graphic designer.

Create a professional medical infographic in SVG format.

Technical specifications:
- Dimensions: 800px width × 1200px height
- Colors: Cyan #06B6D4, Green #10B981, White #FFFFFF
- Text: English only

Create a professional medical design including:
1. Bold title at top
2. Patient information
3. Symptoms with icons
4. Diagnosis in colored frame
5. Treatment recommendations
6. Required medical tests

Return only SVG code without explanation."""

        patient_name = patient_data.get('name', 'Patient')
        symptoms_text = ", ".join(patient_data.get('symptoms', [])[:5])
        diagnosis = patient_data.get('diagnosis', 'Requires medical evaluation')
        recommendations = patient_data.get('recommendations', [])[:4]
        tests = patient_data.get('tests', [])[:4]
        
        user_prompt = f"""Medical Data:

Patient: {patient_name}
Symptoms: {symptoms_text}
Diagnosis: {diagnosis}
Recommendations: {', '.join(recommendations)}
Tests: {', '.join(tests)}

{f"Additional notes: {custom_prompt}" if custom_prompt else ""}

Create a professional medical infographic in English only."""
    
    message = client.messages.create(
        model="claude-sonnet-4-20250514",
        max_tokens=4000,
        temperature=0.3,
        system=system_prompt,
        messages=[{"role": "user", "content": user_prompt}]
    )
    
    svg_content = message.content[0].text
    
    if "```svg" in svg_content:
        svg_content = svg_content.split("```svg")[1].split("```")[0].strip()
    elif "```xml" in svg_content:
        svg_content = svg_content.split("```xml")[1].split("```")[0].strip()
    elif "```" in svg_content:
        svg_content = svg_content.split("```")[1].split("```")[0].strip()
    
    return svg_content

# ==================== PPTX GENERATOR ====================

def create_ultimate_medical_presentation(patient_data: PatientData, consultation_id: int, language: str) -> Presentation:
    """Create ultimate medical presentation matching your exact design!"""
    
    is_arabic = language == "ar"
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    
    # SLIDE 1: COVER
    slide = add_blank_slide(prs)
    add_header_banner(slide, is_arabic)
    
    title_text = "ملخص الحالة الطبية" if is_arabic else "Medical Case Summary"
    add_title(slide, title_text, is_arabic)
    
    patient_text = f"المريض: {patient_data.name}" if is_arabic else f"Patient: {patient_data.name}"
    add_subtitle(slide, patient_text, is_arabic)
    
    # SLIDE 2: SYMPTOMS
    slide = add_blank_slide(prs)
    add_header_banner(slide, is_arabic)
    
    section_title = "الأعراض المقدمة" if is_arabic else "Presented Symptoms"
    add_section_card(slide, title=section_title, items=patient_data.symptoms[:6], 
                     color=BrandColors.CYAN_ACCENT, is_arabic=is_arabic, 
                     left_position=Inches(5.5) if is_arabic else Inches(0.5))
    
    # SLIDE 3: DIAGNOSIS
    slide = add_blank_slide(prs)
    add_header_banner(slide, is_arabic)
    
    diag_title = "التشخيص الأولي" if is_arabic else "Initial Diagnosis"
    add_diagnosis_card(slide, title=diag_title, diagnosis=patient_data.diagnosis, 
                       urgency=patient_data.urgency or "medium", is_arabic=is_arabic)
    
    # SLIDE 4: RECOMMENDATIONS
    slide = add_blank_slide(prs)
    add_header_banner(slide, is_arabic)
    
    rec_title = "التوصيات العلاجية" if is_arabic else "Treatment Recommendations"
    add_section_card(slide, title=rec_title, items=patient_data.recommendations[:6], 
                     color=BrandColors.CYAN_ACCENT, is_arabic=is_arabic, 
                     left_position=Inches(0.5) if is_arabic else Inches(5.5))
    
    # SLIDE 5: TESTS
    slide = add_blank_slide(prs)
    add_header_banner(slide, is_arabic)
    
    test_title = "الفحوصات المقترحة" if is_arabic else "Recommended Tests"
    add_section_card(slide, title=test_title, items=patient_data.tests[:6], 
                     color=BrandColors.GREEN_SUCCESS, is_arabic=is_arabic, 
                     left_position=Inches(2.5))
    
    # SLIDE 6: FOOTER
    slide = add_blank_slide(prs)
    add_header_banner(slide, is_arabic)
    
    footer_text = "مستشارك الطبي الذكي | Smart Medical Consultant\nPremium Healthcare Technology"
    add_footer(slide, footer_text)
    
    return prs

def add_blank_slide(prs: Presentation):
    blank_layout = prs.slide_layouts[6]
    return prs.slides.add_slide(blank_layout)

def add_header_banner(slide, is_arabic: bool):
    banner = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1.25))
    banner.fill.solid()
    banner.fill.fore_color.rgb = BrandColors.CYAN_HEADER
    banner.line.fill.background()

def add_title(slide, text: str, is_arabic: bool):
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    text_frame = title_box.text_frame
    text_frame.text = text
    p = text_frame.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = BrandColors.WHITE
    p.alignment = PP_ALIGN.RIGHT if is_arabic else PP_ALIGN.CENTER

def add_subtitle(slide, text: str, is_arabic: bool):
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.85), Inches(9), Inches(0.3))
    text_frame = subtitle_box.text_frame
    text_frame.text = text
    p = text_frame.paragraphs[0]
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = BrandColors.WHITE
    p.alignment = PP_ALIGN.RIGHT if is_arabic else PP_ALIGN.CENTER

def add_section_card(slide, title: str, items: List[str], color: RGBColor, is_arabic: bool, left_position: float):
    card = slide.shapes.add_shape(1, left_position, Inches(1.5), Inches(4), Inches(3.5))
    card.fill.solid()
    card.fill.fore_color.rgb = BrandColors.WHITE
    card.line.fill.background()
    
    border = slide.shapes.add_shape(1, left_position - Inches(0.05) if is_arabic else left_position + Inches(4),
                                    Inches(1.5), Inches(0.05), Inches(3.5))
    border.fill.solid()
    border.fill.fore_color.rgb = color
    border.line.fill.background()
    
    title_box = slide.shapes.add_textbox(left_position + Inches(0.2), Inches(1.6), Inches(3.6), Inches(0.4))
    text_frame = title_box.text_frame
    text_frame.text = title
    p = text_frame.paragraphs[0]
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = color
    p.alignment = PP_ALIGN.RIGHT if is_arabic else PP_ALIGN.LEFT
    
    items_box = slide.shapes.add_textbox(left_position + Inches(0.3), Inches(2.1), Inches(3.4), Inches(2.8))
    text_frame = items_box.text_frame
    text_frame.word_wrap = True
    
    for item in items[:6]:
        p = text_frame.add_paragraph()
        p.text = f"• {item}" if not is_arabic else f"{item} •"
        p.font.size = Pt(11)
        p.font.color.rgb = BrandColors.TEXT_PRIMARY
        p.alignment = PP_ALIGN.RIGHT if is_arabic else PP_ALIGN.LEFT
        p.space_before = Pt(6)

def add_diagnosis_card(slide, title: str, diagnosis: str, urgency: str, is_arabic: bool):
    card = slide.shapes.add_shape(1, Inches(2.5), Inches(1.8), Inches(5), Inches(2.5))
    card.fill.solid()
    card.fill.fore_color.rgb = BrandColors.GREEN_LIGHT
    card.line.color.rgb = BrandColors.GREEN_SUCCESS
    card.line.width = Pt(2)
    
    title_box = slide.shapes.add_textbox(Inches(2.7), Inches(2), Inches(4.6), Inches(0.4))
    text_frame = title_box.text_frame
    text_frame.text = title
    p = text_frame.paragraphs[0]
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = BrandColors.GREEN_DARK
    p.alignment = PP_ALIGN.RIGHT if is_arabic else PP_ALIGN.LEFT
    
    diag_box = slide.shapes.add_textbox(Inches(2.7), Inches(2.6), Inches(4.6), Inches(1.4))
    text_frame = diag_box.text_frame
    text_frame.text = diagnosis
    text_frame.word_wrap = True
    p = text_frame.paragraphs[0]
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = BrandColors.GREEN_DARK
    p.alignment = PP_ALIGN.RIGHT if is_arabic else PP_ALIGN.LEFT

def add_footer(slide, text: str):
    footer_box = slide.shapes.add_textbox(Inches(0), Inches(4.8), Inches(10), Inches(0.8))
    text_frame = footer_box.text_frame
    text_frame.text = text
    p = text_frame.paragraphs[0]
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = BrandColors.CYAN_ACCENT
    p.alignment = PP_ALIGN.CENTER

# ==================== RUN SERVER ====================

if __name__ == "__main__":
    import uvicorn
    print("""
    ╔══════════════════════════════════════════════════════════════╗
    ║  SMART MEDICAL CONSULTANT API - ULTIMATE EDITION v2.0        ║
    ║  مستشارك الطبي الذكي                                         ║
    ║                                                              ║
    ║  Status: 🚀 READY TO LAUNCH                                  ║
    ║  Features: ✅ PPTX Generation | ✅ Infographics              ║
    ║  Quality: 💎 WORLD-CLASS                                     ║
    ╚══════════════════════════════════════════════════════════════╝
    """)
    uvicorn.run(app, host="0.0.0.0", port=8000, log_level="info")